import streamlit as st
import base64
import pptx
from pptx.util import Inches, Pt
import os
import time
from dotenv import load_dotenv
import google.generativeai as genai

# Load environment variables
load_dotenv()

# Configure the Generative AI client using your API key
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

# Custom formatting options for the presentation
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

def generate_slide_titles(topic):
    """Generates slide titles for a given topic using the Generative AI client."""
    prompt = f"Generate 5 slide titles for the topic '{topic}'."
    response = genai.generate_text(prompt=prompt)
    titles = response.result.split("\n")
    return [title.strip() for title in titles if title.strip()]

def generate_slide_content(slide_title):
    """Generates detailed slide content based on a slide title."""
    prompt = f"Generate detailed content for a slide titled '{slide_title}'."
    response = genai.generate_text(prompt=prompt)
    return response.result.strip()

def create_presentation(topic, slide_titles, slide_contents):
    """Creates a PowerPoint presentation from slide titles and their content."""
    prs = pptx.Presentation()
    # Create a title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    # Create a slide for each title-content pair
    for title, content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.shapes.placeholders[1].text = content

        # Set custom font sizes
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

    # Save the presentation in a directory
    output_dir = "generated_ppt"
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    filename = os.path.join(output_dir, f"{topic}_presentation.pptx")
    prs.save(filename)
    return filename

def get_ppt_download_link(filepath):
    """Generates an HTML download link for the PPT file."""
    with open(filepath, "rb") as f:
        ppt_data = f.read()
    b64 = base64.b64encode(ppt_data).decode()
    return (
        f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" '
        f'download="{os.path.basename(filepath)}">Download Presentation</a>'
    )

def main():
    st.title("PowerPoint Presentation Generator using Google Gen AI (Gemini)")
    topic = st.text_input("Enter the topic for your presentation:")
    if st.button("Generate Presentation") and topic:
        st.info("Generating presentation, please wait...")
        start_time = time.time()

        slide_titles = generate_slide_titles(topic)
        slide_contents = [generate_slide_content(title) for title in slide_titles]
        ppt_file = create_presentation(topic, slide_titles, slide_contents)
        elapsed = time.time() - start_time
        st.success(f"Presentation generated in {elapsed:.2f} seconds!")
        st.markdown(get_ppt_download_link(ppt_file), unsafe_allow_html=True)

if _name_ == "_main_":
    main()